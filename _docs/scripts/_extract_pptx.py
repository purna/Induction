"""Extract text from PowerPoint decks in _induction/{y1,y2}/.

Output: /tmp/induction_txt/{year}/{slug}.txt
Each file contains:
  SLIDE 1: <title>
    bullet
    bullet
  NOTES:
    speaker note

This script is read-only on _induction/.
"""
import os
import re
import sys
from pathlib import Path

from pptx import Presentation


REPO = Path(__file__).resolve().parent
SRC = REPO / "_induction"
OUT = Path("/tmp/induction_txt")

YEAR_DIRS = {
    "y1": SRC / "y1",
    "y2": SRC / "y2",
}


def slide_bullets(slide):
    lines = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        for para in tf.paragraphs:
            txt = "".join(run.text for run in para.runs).strip()
            if not txt:
                continue
            level = para.level or 0
            indent = "  " * level
            lines.append(f"{indent}- {txt}")
    return lines


def slide_notes(slide):
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            return notes
    return ""


def safe_slug(name):
    s = re.sub(r"[^a-zA-Z0-9]+", "_", name).strip("_").lower()
    return s


def main():
    for year, src_dir in YEAR_DIRS.items():
        if not src_dir.exists():
            print(f"missing: {src_dir}", file=sys.stderr)
            continue
        out_dir = OUT / year
        out_dir.mkdir(parents=True, exist_ok=True)
        for path in sorted(src_dir.glob("*.pptx")):
            prs = Presentation(str(path))
            lines = []
            for i, slide in enumerate(prs.slides, start=1):
                title = ""
                try:
                    if slide.shapes.title and slide.shapes.title.has_text_frame:
                        title = slide.shapes.title.text_frame.text.strip()
                except Exception:
                    pass
                lines.append(f"SLIDE {i}: {title}")
                for b in slide_bullets(slide):
                    lines.append(f"  {b}")
                notes = slide_notes(slide)
                if notes:
                    lines.append("  NOTES:")
                    for nl in notes.splitlines():
                        lines.append(f"    {nl}")
                lines.append("")
            slug = safe_slug(path.stem)
            out_path = out_dir / f"{slug}.txt"
            out_path.write_text("\n".join(lines), encoding="utf-8")
            print(f"wrote {out_path}")


if __name__ == "__main__":
    main()